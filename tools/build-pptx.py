#!/usr/bin/env python3
"""
build-pptx.py — Generate the C++ container optimization tutorial deck.

Produces presentation/cpp-container-tutorial.pptx (3-hour cut, ~110
slides with full speaker notes) from the section content in _docs/
plus the SVG diagrams in diagrams/.

Design tokens borrowed from patterncatalyst/quarkus-optimization's
deck for visual continuity with the author's other talks.

Usage:
    python3 tools/build-pptx.py

Output:
    presentation/cpp-container-tutorial.pptx
    presentation/build-notes.md          # generation log

Prerequisites:
    pip install python-pptx
    LibreOffice (soffice) for SVG→PNG conversion (done by tools/convert-diagrams.sh first)
"""
from __future__ import annotations

import json
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================================
# Design tokens — extracted from the Quarkus deck
# ============================================================================

class C:
    """Color palette — dark navy + accent system."""
    # Backgrounds
    BG_TITLE_DARK    = RGBColor(0x1A, 0x2B, 0x3C)   # title slide bg
    BG_SECTION_DARK  = RGBColor(0x0A, 0x16, 0x28)   # section divider bg
    BG_HEADER_BAR    = RGBColor(0x1E, 0x6F, 0xC8)   # content slide header bar
    BG_CODE_BLOCK    = RGBColor(0x0D, 0x21, 0x37)   # code block bg
    BG_CARD_LIGHT    = RGBColor(0xEC, 0xF0, 0xF1)   # neutral light card
    BG_CARD_SOFT     = RGBColor(0xE0, 0xE8, 0xF0)   # soft blue card

    # Text
    TEXT_DARK        = RGBColor(0x1A, 0x2B, 0x3C)   # body text on light bg
    TEXT_LIGHT       = RGBColor(0xEC, 0xF0, 0xF1)   # body text on dark bg
    TEXT_WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # pure white
    TEXT_MUTED       = RGBColor(0x90, 0xA4, 0xAE)   # muted secondary

    # Accents
    ACCENT_CYAN      = RGBColor(0x00, 0xBC, 0xD4)   # primary accent
    ACCENT_BLUE      = RGBColor(0x1E, 0x6F, 0xC8)   # blue accent
    ACCENT_GREEN     = RGBColor(0x27, 0xAE, 0x60)   # "after" / positive
    ACCENT_RED       = RGBColor(0xE8, 0x48, 0x55)   # "before" / negative
    ACCENT_ORANGE    = RGBColor(0xF5, 0xA6, 0x23)   # warnings
    ACCENT_PURPLE    = RGBColor(0x9B, 0x59, 0xB6)   # tertiary

    # Demo cue color (distinctive — matches the "DEMO" pill style)
    BG_DEMO_PILL     = RGBColor(0x27, 0xAE, 0x60)


class F:
    """Font sizes."""
    TITLE_HUGE     = Pt(44)   # title slide
    TITLE_LARGE    = Pt(36)   # section divider
    TITLE_MED      = Pt(28)   # slide title
    SECTION_LABEL  = Pt(14)   # "Section 02" pill
    BODY_LARGE     = Pt(20)   # large body
    BODY           = Pt(16)   # standard body
    BODY_SMALL     = Pt(14)   # small body
    CAPTION        = Pt(11)   # captions, footer
    CODE           = Pt(12)   # code blocks
    CODE_SMALL     = Pt(11)   # long code blocks
    STAT_HUGE      = Pt(72)   # big-number callouts
    STAT_LARGE     = Pt(48)


class FontFam:
    HEADER = "Calibri"
    BODY   = "Calibri"
    CODE   = "Consolas"


# Slide dimensions: 16:9, 13.333" x 7.5"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN  = Inches(0.5)


# ============================================================================
# Project paths
# ============================================================================

ROOT = Path(__file__).resolve().parent.parent
DIAGRAMS_DIR = Path("/tmp/diagrams-png")   # converted PNGs
OUT_PATH = ROOT / "presentation" / "cpp-container-tutorial.pptx"
NOTES_PATH = ROOT / "presentation" / "build-notes.md"
BIBLIO_URL = "https://patterncatalyst.github.io/cpp-container-optimization-tutorial/bibliography/"
REPO_URL = "https://github.com/patterncatalyst/cpp-container-optimization-tutorial"


# ============================================================================
# Helpers — slide building primitives
# ============================================================================

def set_slide_bg(slide, color: RGBColor) -> None:
    """Set a solid background fill on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor,
             line_color: Optional[RGBColor] = None):
    """Add a rectangle shape. Returns the shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    return shape


def add_pill(slide, left, top, width, height, fill_color: RGBColor):
    """Add a rounded-rectangle 'pill' shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Tight corner radius
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    return shape


def add_text(slide, left, top, width, height, text: str,
             size: Pt = F.BODY,
             color: RGBColor = C.TEXT_DARK,
             bold: bool = False,
             italic: bool = False,
             align: PP_ALIGN = PP_ALIGN.LEFT,
             anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
             font: str = FontFam.BODY):
    """Add a text box with single paragraph of styled text."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    return tb


def add_multi_text(slide, left, top, width, height,
                   paragraphs: list[dict],
                   anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
                   line_spacing: float = 1.15):
    """Add a text box with multiple styled paragraphs.

    Each paragraph dict: {text, size, color, bold, italic, align, font, bullet}
    """
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get("align", PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        if para.get("space_before"):
            p.space_before = para["space_before"]

        text = para["text"]
        # Bullet handling: prefix with • if requested
        if para.get("bullet"):
            text = "•  " + text

        run = p.add_run()
        run.text = text
        run.font.size = para.get("size", F.BODY)
        run.font.bold = para.get("bold", False)
        run.font.italic = para.get("italic", False)
        run.font.color.rgb = para.get("color", C.TEXT_DARK)
        run.font.name = para.get("font", FontFam.BODY)
    return tb


def add_header_bar(slide, section_label: str, page_no: int, total_pages: int,
                   section_title: str = ""):
    """Add the standard navy header bar across the top of content slides."""
    # Full-width header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.7), C.BG_HEADER_BAR)

    # Section label on left
    add_text(slide, Inches(0.4), Inches(0.15), Inches(7), Inches(0.4),
             section_title or section_label,
             size=F.TITLE_MED, color=C.TEXT_WHITE,
             bold=True, align=PP_ALIGN.LEFT,
             font=FontFam.HEADER)

    # Section pill on right
    pill_w = Inches(1.6)
    pill = add_pill(slide,
                    SLIDE_W - pill_w - Inches(0.3),
                    Inches(0.18),
                    pill_w, Inches(0.35),
                    C.BG_TITLE_DARK)
    add_text(slide, SLIDE_W - pill_w - Inches(0.3), Inches(0.18),
             pill_w, Inches(0.35),
             section_label,
             size=F.SECTION_LABEL, color=C.TEXT_WHITE,
             bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FontFam.HEADER)


def add_footer(slide, section_title: str, page_no: int, total_pages: int):
    """Add the standard footer bar."""
    footer_top = SLIDE_H - Inches(0.35)
    add_text(slide, Inches(0.4), footer_top,
             Inches(9), Inches(0.3),
             f"Optimizing Modern C++ with Containers   |   {section_title}",
             size=F.CAPTION, color=C.TEXT_MUTED,
             align=PP_ALIGN.LEFT, font=FontFam.BODY)
    add_text(slide, SLIDE_W - Inches(1.5), footer_top,
             Inches(1.1), Inches(0.3),
             f"{page_no} / {total_pages}",
             size=F.CAPTION, color=C.TEXT_MUTED,
             align=PP_ALIGN.RIGHT, font=FontFam.BODY)


def set_notes(slide, speaker_text: str) -> None:
    """Add full speaker notes to a slide."""
    notes = slide.notes_slide
    tf = notes.notes_text_frame
    # Clear existing
    tf.clear()
    # Split into paragraphs by double-newline
    paragraphs = [p.strip() for p in speaker_text.split("\n\n") if p.strip()]
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = para
        run.font.size = Pt(12)


# ============================================================================
# Slide builders — title, section divider, content, demo cue, code, diagram
# ============================================================================

def build_title_slide(prs, total_pages: int) -> None:
    """Slide 1 — title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, C.BG_TITLE_DARK)

    # Three accent dots
    dot_y = Inches(1.5)
    for i, color in enumerate([C.ACCENT_CYAN, C.ACCENT_ORANGE, C.ACCENT_RED]):
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0 + i * 0.35), dot_y, Inches(0.22), Inches(0.22)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()

    # Eyebrow
    add_text(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(0.5),
             "OPTIMIZING MODERN C++",
             size=Pt(20), color=C.ACCENT_CYAN,
             bold=True, font=FontFam.HEADER)

    # Main title
    add_text(slide, Inches(1.0), Inches(2.5), Inches(11), Inches(1.5),
             "C++20/23 Performance Under\nContainer Constraints",
             size=F.TITLE_HUGE, color=C.TEXT_WHITE,
             bold=True, font=FontFam.HEADER)

    # Tech stack subtitle
    add_text(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(0.5),
             "Fedora 44  •  Podman 5.x  •  GCC 14 / Clang 18  •  C++23  "
             "•  io_uring  •  cgroups v2  •  LGTM stack",
             size=F.BODY, color=C.TEXT_LIGHT, font=FontFam.BODY)

    # Source references
    add_text(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(0.4),
             "Based on: Andrist & Sehr, Iglberger, Enberg, Ghosh",
             size=F.BODY_SMALL, color=C.TEXT_MUTED,
             italic=True, font=FontFam.BODY)

    # Talk metadata
    add_text(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(0.4),
             "180-minute deep dive  •  7 runnable demos  •  17 sections",
             size=F.BODY_SMALL, color=C.ACCENT_CYAN, font=FontFam.BODY)

    # Repo
    add_text(slide, Inches(1.0), Inches(6.5), Inches(11), Inches(0.4),
             "github.com/patterncatalyst/cpp-container-optimization-tutorial",
             size=F.CAPTION, color=C.TEXT_MUTED, font=FontFam.CODE)

    set_notes(slide,
        "Welcome. This is a three-hour deep dive on optimizing modern C++ — "
        "the C++20 and C++23 we ship today — inside the operational reality "
        "of OCI containers on Linux. Not theory: every claim has a runnable "
        "demo on the public repo.\n\n"
        "Three things you should know going in. First, this is a working "
        "Fedora 44 plus Podman 5 lab; the only host you need is the one "
        "you're sitting at. Second, the seven demos shipped with the "
        "tutorial run end-to-end with one shell script each — there is no "
        "external test environment. Third, the design constraints are "
        "production-realistic: cgroup limits, rootless containers, the "
        "Grafana LGTM observability stack, real CVE-aware base images.\n\n"
        "If you're here because someone asked you to make a C++ service "
        "faster and the advice you found assumed bare metal — this talk is "
        "the bridge. We'll walk from compile-time decisions (LTO, PGO, "
        "constexpr) through data-structure choice (the flat_map vs "
        "unordered_map question at scale), through allocator selection "
        "(PMR vs mimalloc vs the default), through I/O strategy (io_uring "
        "direct vs Asio wrapped, async gRPC), through observability and "
        "isolation, and end at reproducibility and ABI hygiene.\n\n"
        "Quick housekeeping: the slides have a section pill in the top-"
        "right showing where we are. The companion site has every code "
        "listing, every metric, and every diagram in long form. The "
        "repository has all seven demos with verified output. Q&A buffer "
        "is at the end; live demos run inline. Let's go.")


def build_agenda_slide(prs, total_pages: int) -> None:
    """Slide 2 — agenda grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.TEXT_WHITE)

    add_header_bar(slide, "Overview", 2, total_pages, "Agenda")
    add_footer(slide, "Agenda", 2, total_pages)

    # Two columns of 8-9 sections each
    sections = [
        ("§1", "Prerequisites", C.ACCENT_RED),
        ("§2", "Introduction & mental model", C.ACCENT_ORANGE),
        ("§3", "RAII & resource discipline", C.ACCENT_BLUE),
        ("§4", "Container strategy: UBI, multi-stage", C.ACCENT_GREEN),
        ("§5", "Compile-time wins: LTO, PGO, constexpr", C.ACCENT_CYAN),
        ("§6", "STL, layout, C++20/23 containers", C.ACCENT_PURPLE),
        ("§7", "Memory: allocators, huge pages, cgroups", C.ACCENT_RED),
        ("§8", "I/O latency: io_uring, async gRPC", C.ACCENT_ORANGE),
        ("§9", "Networking & kernel parameters", C.ACCENT_BLUE),
        ("§10", "Observability: OTel, perf, eBPF", C.ACCENT_GREEN),
        ("§11", "Noisy-neighbor isolation: cgroups, NUMA", C.ACCENT_CYAN),
        ("§12", "Static analysis & debugging", C.ACCENT_PURPLE),
        ("§13", "Reproducibility & ABI", C.ACCENT_RED),
        ("§14", "Pitfalls", C.ACCENT_ORANGE),
        ("§15", "Where to go next + bibliography", C.ACCENT_BLUE),
        ("§16", "Appendix — Conan + UBI 9 perl", C.TEXT_MUTED),
    ]
    # 2-column grid: 8 per column
    cols = 2
    rows_per_col = 8
    col_width = Inches(6.0)
    row_height = Inches(0.62)
    grid_left = Inches(0.5)
    grid_top = Inches(1.1)

    for i, (num, title, color) in enumerate(sections):
        col = i // rows_per_col
        row = i % rows_per_col
        x = grid_left + col * (col_width + Inches(0.3))
        y = grid_top + row * row_height

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y, Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.fill.background()
        add_text(slide, x, y, Inches(0.5), Inches(0.5),
                 num, size=Pt(13), color=C.TEXT_WHITE,
                 bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font=FontFam.HEADER)

        # Title
        add_text(slide, x + Inches(0.65), y + Inches(0.05),
                 col_width - Inches(0.7), Inches(0.4),
                 title,
                 size=F.BODY, color=C.TEXT_DARK,
                 anchor=MSO_ANCHOR.MIDDLE, font=FontFam.BODY)

    set_notes(slide,
        "Sixteen numbered sections plus an appendix. We'll do every one "
        "of them in three hours. The reading time on the companion site is "
        "longer — closer to ten or twelve hours if you actually run every "
        "demo — but the talking flow is tight.\n\n"
        "Here's the shape. Sections 1, 2, and 3 are setup and mental model: "
        "what we're optimizing, where the levers live, and the RAII "
        "discipline that the whole tutorial assumes. Sections 4 and 5 are "
        "compile-time work — image strategy and LTO/PGO — with Demo 1 as "
        "the worked example. Section 6 is the data-structure question; "
        "Demo 2. Section 7 is allocator strategy; Demo 6. Sections 8 and 9 "
        "are I/O and networking; Demo 3. Section 10 is observability — Demo "
        "4 is the Grafana LGTM stack. Section 11 is noisy-neighbor isolation; "
        "Demo 5. Sections 12 and 13 are the engineering hygiene track: "
        "analysis, debugging, reproducibility, ABI; Demo 7. Section 14 is "
        "pitfalls — the runbook. Section 15 is reading pointers. Section 16 "
        "is an appendix you only need if you're doing Conan from-source on "
        "UBI 9.\n\n"
        "Every demo runs with one shell script. Every result you'll see "
        "today was reproduced on Fedora 44, Podman 5.x, GCC 14, rootless. "
        "If anything I say doesn't match what your machine does, open an "
        "issue with your uname-r and lscpu output.")


def build_section_divider(prs, section_num: int, title: str,
                          tagline: str, page_no: int, total_pages: int,
                          notes: str) -> None:
    """Section divider slide — dark navy bg, big number + title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.BG_SECTION_DARK)

    # Big section number
    add_text(slide, Inches(0.5), Inches(2.0), Inches(4), Inches(2.5),
             f"§{section_num}",
             size=Pt(180), color=C.ACCENT_CYAN,
             bold=True, font=FontFam.HEADER)

    # Section title
    add_text(slide, Inches(4.5), Inches(2.3), Inches(8.5), Inches(2),
             title,
             size=F.TITLE_LARGE, color=C.TEXT_WHITE,
             bold=True, font=FontFam.HEADER)

    # Tagline
    add_text(slide, Inches(4.5), Inches(3.8), Inches(8.5), Inches(2),
             tagline,
             size=F.BODY_LARGE, color=C.TEXT_MUTED,
             italic=True, font=FontFam.BODY)

    # Page indicator
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.5),
             Inches(1.0), Inches(0.3),
             f"{page_no} / {total_pages}",
             size=F.CAPTION, color=C.TEXT_MUTED,
             align=PP_ALIGN.RIGHT, font=FontFam.BODY)

    set_notes(slide, notes)


def build_content_slide(prs, section_label: str, section_title: str,
                        title: str,
                        body_paragraphs: list[dict],
                        page_no: int, total_pages: int,
                        notes: str,
                        right_content: Optional[dict] = None) -> None:
    """Standard two-column or full-width content slide.

    body_paragraphs: list of dicts for add_multi_text
    right_content: optional dict {type: 'code'|'image'|'card', ...}
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.TEXT_WHITE)

    add_header_bar(slide, section_label, page_no, total_pages, section_title)
    add_footer(slide, section_title, page_no, total_pages)

    # Slide title
    add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.6),
             title,
             size=F.TITLE_MED, color=C.TEXT_DARK,
             bold=True, font=FontFam.HEADER)

    # Body region
    if right_content is None:
        # Full-width body
        body_left = Inches(0.5)
        body_width = Inches(12.3)
    else:
        # Half-width body, right column for content
        body_left = Inches(0.5)
        body_width = Inches(6.0)

    add_multi_text(slide, body_left, Inches(1.7),
                   body_width, Inches(5.0),
                   body_paragraphs)

    # Right content
    if right_content:
        rc_left = Inches(6.8)
        rc_top = Inches(1.7)
        rc_width = Inches(6.0)
        rc_height = Inches(5.0)
        rc_type = right_content.get("type")
        if rc_type == "code":
            add_code_block(slide, rc_left, rc_top, rc_width, rc_height,
                           right_content["code"],
                           lang=right_content.get("lang", "cpp"))
        elif rc_type == "image":
            img_path = right_content["path"]
            if Path(img_path).exists():
                slide.shapes.add_picture(img_path, rc_left, rc_top,
                                         width=rc_width)
        elif rc_type == "card":
            add_card(slide, rc_left, rc_top, rc_width, rc_height,
                     right_content["title"], right_content["body"],
                     bg=right_content.get("bg", C.BG_CARD_LIGHT))
        elif rc_type == "stat":
            add_stat_callout(slide, rc_left, rc_top, rc_width, rc_height,
                             right_content["number"],
                             right_content["label"],
                             color=right_content.get("color", C.ACCENT_CYAN))

    set_notes(slide, notes)


def add_code_block(slide, left, top, width, height,
                   code: str, lang: str = "cpp"):
    """Add a dark-navy code block. Auto-shrinks for long content."""
    add_rect(slide, left, top, width, height, C.BG_CODE_BLOCK)
    tb = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), height - Inches(0.2)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.05)
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = code.split("\n")
    max_line = max(len(l) for l in lines) if lines else 0
    n_lines = len(lines)
    # Pick a font size that should fit. Roughly: 11pt fits ~30 lines in a
    # 5" tall box. 12pt fits ~26 lines. 10pt fits ~33.
    if n_lines > 28 or max_line > 70:
        size = Pt(10)
    elif n_lines > 22 or max_line > 60:
        size = F.CODE_SMALL
    else:
        size = F.CODE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line if line else " "
        run.font.size = size
        run.font.name = FontFam.CODE
        run.font.color.rgb = C.TEXT_LIGHT


def add_card(slide, left, top, width, height,
             title: str, body: str,
             bg: RGBColor = C.BG_CARD_LIGHT,
             title_color: RGBColor = C.TEXT_DARK,
             body_color: RGBColor = C.TEXT_DARK):
    """Add a titled card box."""
    add_rect(slide, left, top, width, height, bg)
    add_text(slide, left + Inches(0.2), top + Inches(0.15),
             width - Inches(0.4), Inches(0.4),
             title, size=F.BODY_LARGE, color=title_color,
             bold=True, font=FontFam.HEADER)
    add_text(slide, left + Inches(0.2), top + Inches(0.7),
             width - Inches(0.4), height - Inches(0.85),
             body, size=F.BODY_SMALL, color=body_color,
             font=FontFam.BODY)


def add_stat_callout(slide, left, top, width, height,
                     number: str, label: str,
                     color: RGBColor = C.ACCENT_CYAN):
    """Big-number callout — number on top, label below."""
    add_rect(slide, left, top, width, height, C.BG_CARD_SOFT)
    # Number
    add_text(slide, left, top + Inches(0.5), width, Inches(2),
             number, size=F.STAT_HUGE, color=color,
             bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FontFam.HEADER)
    # Label
    add_text(slide, left + Inches(0.2), top + Inches(2.8),
             width - Inches(0.4), Inches(1.5),
             label, size=F.BODY, color=C.TEXT_DARK,
             align=PP_ALIGN.CENTER, font=FontFam.BODY)


def build_diagram_slide(prs, section_label: str, section_title: str,
                        title: str,
                        diagram_path: str,
                        caption: str,
                        page_no: int, total_pages: int,
                        notes: str) -> None:
    """A slide whose primary content is a single diagram."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.TEXT_WHITE)

    add_header_bar(slide, section_label, page_no, total_pages, section_title)
    add_footer(slide, section_title, page_no, total_pages)

    # Title
    add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.5),
             title,
             size=F.TITLE_MED, color=C.TEXT_DARK,
             bold=True, font=FontFam.HEADER)

    # Diagram — size by aspect ratio so it fits the available area.
    # Available: 12.3" wide x 5.0" tall (between title and caption)
    avail_w = Inches(12.3)
    avail_h = Inches(5.0)
    if Path(diagram_path).exists():
        try:
            from PIL import Image
            with Image.open(diagram_path) as img:
                px_w, px_h = img.size
            aspect = px_w / px_h
            # Fit to whichever dimension hits the cap first
            if aspect >= avail_w / avail_h:
                # Wider than tall → constrain by width
                draw_w = avail_w
                draw_h = int(avail_w / aspect)
            else:
                # Taller than wide → constrain by height
                draw_h = avail_h
                draw_w = int(avail_h * aspect)
            # Center horizontally
            left = Inches(0.5) + (avail_w - draw_w) / 2
            top = Inches(1.6) + (avail_h - draw_h) / 2
            slide.shapes.add_picture(diagram_path, left, top,
                                     width=draw_w, height=draw_h)
        except Exception:
            slide.shapes.add_picture(diagram_path, Inches(1.0),
                                     Inches(1.6), width=Inches(11.3))
    else:
        add_text(slide, Inches(1.0), Inches(3.0), Inches(11.3), Inches(2),
                 f"[Diagram not found: {diagram_path}]",
                 size=F.BODY, color=C.ACCENT_RED, align=PP_ALIGN.CENTER)

    # Caption
    add_text(slide, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.35),
             caption,
             size=F.CAPTION, color=C.TEXT_MUTED,
             italic=True, align=PP_ALIGN.CENTER, font=FontFam.BODY)

    set_notes(slide, notes)


def build_demo_cue(prs, section_label: str, section_title: str,
                   demo_num: int, demo_name: str,
                   demo_command: str, demo_description: str,
                   demo_url: str,
                   page_no: int, total_pages: int,
                   notes: str) -> None:
    """A demo cue slide — green DEMO pill, command, description."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.BG_SECTION_DARK)

    # "DEMO N" badge
    pill_w = Inches(2.5)
    add_pill(slide, Inches(0.6), Inches(0.6), pill_w, Inches(0.65),
             C.BG_DEMO_PILL)
    add_text(slide, Inches(0.6), Inches(0.6), pill_w, Inches(0.65),
             f"▶ DEMO {demo_num:02d}",
             size=Pt(22), color=C.TEXT_WHITE,
             bold=True, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE, font=FontFam.HEADER)

    # Demo name
    add_text(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
             demo_name,
             size=F.TITLE_LARGE, color=C.TEXT_WHITE,
             bold=True, font=FontFam.HEADER)

    # Description
    add_text(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(1.5),
             demo_description,
             size=F.BODY_LARGE, color=C.TEXT_LIGHT, font=FontFam.BODY)

    # Command — code block style
    cmd_top = Inches(4.7)
    add_rect(slide, Inches(0.6), cmd_top, Inches(12.1), Inches(1.2),
             C.BG_CODE_BLOCK)
    add_text(slide, Inches(0.8), cmd_top + Inches(0.1),
             Inches(11.8), Inches(0.5),
             "$ " + demo_command,
             size=Pt(18), color=C.ACCENT_CYAN,
             font=FontFam.CODE)
    add_text(slide, Inches(0.8), cmd_top + Inches(0.65),
             Inches(11.8), Inches(0.4),
             f"  source on site: {demo_url}",
             size=F.CAPTION, color=C.TEXT_MUTED, font=FontFam.CODE)

    # Page number
    add_text(slide, SLIDE_W - Inches(1.5), SLIDE_H - Inches(0.5),
             Inches(1.0), Inches(0.3),
             f"{page_no} / {total_pages}",
             size=F.CAPTION, color=C.TEXT_MUTED,
             align=PP_ALIGN.RIGHT, font=FontFam.BODY)

    set_notes(slide, notes)


def build_closing_slide(prs, total_pages: int) -> None:
    """Final slide — thanks + repo + bibliography."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.BG_TITLE_DARK)

    add_text(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
             "Thank you.",
             size=Pt(64), color=C.TEXT_WHITE,
             bold=True, align=PP_ALIGN.CENTER,
             font=FontFam.HEADER)
    add_text(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.6),
             "Every claim has a runnable demo on the repo.",
             size=F.BODY_LARGE, color=C.ACCENT_CYAN,
             italic=True, align=PP_ALIGN.CENTER,
             font=FontFam.BODY)

    # Three callouts: tutorial, demos, bibliography
    box_top = Inches(4.3)
    box_w = Inches(3.9)
    box_h = Inches(2.0)
    gap = Inches(0.2)
    total_w = box_w * 3 + gap * 2
    box_left_start = (SLIDE_W - total_w) / 2

    callouts = [
        ("Tutorial site", "patterncatalyst.github.io/cpp-container-optimization-tutorial",
         C.ACCENT_CYAN),
        ("All 7 demos", "github.com/.../examples", C.ACCENT_GREEN),
        ("Bibliography", "/bibliography/ — Andrist & Sehr, Iglberger,\nEnberg, Ghosh",
         C.ACCENT_ORANGE),
    ]
    for i, (title, body, color) in enumerate(callouts):
        x = box_left_start + i * (box_w + gap)
        add_rect(slide, x, box_top, box_w, box_h, C.BG_SECTION_DARK,
                 line_color=color)
        add_text(slide, x + Inches(0.2), box_top + Inches(0.2),
                 box_w - Inches(0.4), Inches(0.4),
                 title, size=F.BODY_LARGE, color=color,
                 bold=True, font=FontFam.HEADER)
        add_text(slide, x + Inches(0.2), box_top + Inches(0.7),
                 box_w - Inches(0.4), box_h - Inches(0.85),
                 body, size=F.BODY_SMALL, color=C.TEXT_LIGHT,
                 font=FontFam.BODY)

    set_notes(slide,
        "We covered seventeen sections, seven demos, the bibliography, "
        "and the reproducibility story. Three places to go from here:\n\n"
        "The tutorial site at patterncatalyst.github.io is the long-form "
        "reference. Every code listing, every metric, every diagram — and "
        "the reconciliation plan, which is the truth source for what's "
        "verified versus what's claimed.\n\n"
        "The GitHub repo has all seven demos. Each one is a self-contained "
        "Podman project; cd in, run ./demo.sh, and you have the numbers. "
        "Fork it, run the demos against your own workload, file an issue "
        "if a measurement doesn't match what your machine does.\n\n"
        "The bibliography page consolidates the four reference books — "
        "Andrist & Sehr, Iglberger, Enberg, Ghosh — with annotations on "
        "what each is strongest at and which sections of this tutorial "
        "draw on each. If you want one to read next, that page tells you "
        "which one and in what order.\n\n"
        "Questions?")


def build_stat_row(prs, section_label: str, section_title: str,
                   title: str, stats: list,
                   page_no: int, total_pages: int,
                   notes: str) -> None:
    """Row of 4 big-number callouts — the 'gap' slide pattern."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, C.TEXT_WHITE)

    add_header_bar(slide, section_label, page_no, total_pages, section_title)
    add_footer(slide, section_title, page_no, total_pages)

    # Slide title
    add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.6),
             title, size=F.TITLE_MED, color=C.TEXT_DARK,
             bold=True, font=FontFam.HEADER)

    # 4 stat callouts side-by-side
    n = len(stats)
    avail_w = Inches(12.3)
    gap = Inches(0.2)
    box_w = (avail_w - gap * (n - 1)) / n
    box_h = Inches(4.5)
    top = Inches(1.7)
    left_start = Inches(0.5)

    for i, (number, label, color) in enumerate(stats):
        x = left_start + i * (box_w + gap)
        add_rect(slide, x, top, box_w, box_h, C.BG_CARD_SOFT)

        # Number (big)
        add_text(slide, x, top + Inches(0.3), box_w, Inches(2.2),
                 number, size=Pt(54), color=color,
                 bold=True, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE, font=FontFam.HEADER)

        # Label below
        add_text(slide, x + Inches(0.15), top + Inches(2.7),
                 box_w - Inches(0.3), box_h - Inches(2.9),
                 label, size=F.BODY_SMALL, color=C.TEXT_DARK,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
                 font=FontFam.BODY)

    set_notes(slide, notes)


# ============================================================================
# Dispatcher — render a slide based on its 'kind' field
# ============================================================================

def render_slide(prs, section, slide_data: dict,
                 page_no: int, total_pages: int) -> None:
    """Dispatch on slide kind."""
    kind = slide_data["kind"]
    label = section["label"]
    section_title = section["title"]

    if kind == "content":
        right = None
        if slide_data.get("code"):
            right = {"type": "code", "code": slide_data["code"]}
        elif slide_data.get("diagram"):
            right = {"type": "image", "path": slide_data["diagram"]}
        build_content_slide(
            prs, label, section_title,
            slide_data["title"],
            slide_data["body"],
            page_no, total_pages,
            slide_data["notes"],
            right_content=right,
        )

    elif kind == "content-code":
        right = {"type": "code", "code": slide_data["code"]}
        build_content_slide(
            prs, label, section_title,
            slide_data["title"],
            slide_data["body"],
            page_no, total_pages,
            slide_data["notes"],
            right_content=right,
        )

    elif kind == "code-content":
        # code goes on left, body on right — flipped layout
        # For simplicity, render as content-code (body left, code right)
        right = {"type": "code", "code": slide_data["code"]}
        build_content_slide(
            prs, label, section_title,
            slide_data["title"],
            slide_data["body"],
            page_no, total_pages,
            slide_data["notes"],
            right_content=right,
        )

    elif kind == "diagram":
        build_diagram_slide(
            prs, label, section_title,
            slide_data["title"],
            slide_data["diagram"],
            slide_data["caption"],
            page_no, total_pages,
            slide_data["notes"],
        )

    elif kind == "stat-row":
        build_stat_row(
            prs, label, section_title,
            slide_data["title"],
            slide_data["stats"],
            page_no, total_pages,
            slide_data["notes"],
        )

    elif kind == "demo-cue":
        build_demo_cue(
            prs, label, section_title,
            slide_data["demo_num"],
            slide_data["demo_name"],
            slide_data["demo_command"],
            slide_data["demo_description"],
            slide_data["demo_url"],
            page_no, total_pages,
            slide_data["notes"],
        )

    else:
        raise ValueError(f"Unknown slide kind: {kind}")


# ============================================================================
# Main — read sections.py, render every slide, save
# ============================================================================

def main():
    # Import section data
    sys.path.insert(0, str(Path(__file__).parent))
    from sections import SECTIONS

    # Compute total page count
    total = 2   # title + agenda
    for section in SECTIONS:
        total += 1   # section divider
        total += len(section["slides"])
    total += 1   # closing

    print(f"Total slides planned: {total}")

    # Build the presentation
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    page = 1
    print(f"  [{page:3d}/{total}] Title slide")
    build_title_slide(prs, total)
    page += 1

    print(f"  [{page:3d}/{total}] Agenda slide")
    build_agenda_slide(prs, total)
    page += 1

    for section in SECTIONS:
        num = section["num"]
        title = section["title"]
        print(f"  [{page:3d}/{total}] §{num} divider — {title}")
        build_section_divider(
            prs, num, title,
            section["tagline"],
            page, total,
            section["divider_notes"],
        )
        page += 1

        for slide_data in section["slides"]:
            kind = slide_data["kind"]
            slide_title = slide_data.get("title") or \
                          slide_data.get("demo_name", "(demo)")
            print(f"  [{page:3d}/{total}] §{num} {kind} — {slide_title[:50]}")
            render_slide(prs, section, slide_data, page, total)
            page += 1

    print(f"  [{page:3d}/{total}] Closing slide")
    build_closing_slide(prs, total)
    page += 1

    # Save
    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    print(f"\nSaved: {OUT_PATH}")
    print(f"Slide count: {page - 1}")

    # Write a brief build-notes log
    NOTES_PATH.write_text(
        f"# Build notes — cpp-container-tutorial.pptx\n\n"
        f"- Total slides: {page - 1}\n"
        f"- Sections rendered: {len(SECTIONS)}\n"
        f"- Diagrams embedded: from {DIAGRAMS_DIR}\n"
        f"- Generator: tools/build-pptx.py\n"
        f"- Content source: tools/sections.py\n"
    )
    print(f"Wrote: {NOTES_PATH}")


if __name__ == "__main__":
    main()
